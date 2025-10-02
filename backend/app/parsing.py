# app/parsing.py
from __future__ import annotations
from typing import Iterable, Iterator, Tuple, List
import os
import io
import re
import zipfile
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fitz  # PyMuPDF
from PIL import Image


# ------------------------------
# Generic helpers
# ------------------------------

def _clean_text(s: str) -> str:
    """Normalize whitespace, remove nulls and tame hyphenated line breaks."""
    if not s:
        return ""
    s = s.replace("\r\n", "\n").replace("\r", "\n")
    s = s.replace("\x00", " ")
    # Join words split by line breaks with hyphen (e.g., "re-\nsolve" -> "resolve")
    s = re.sub(r"-\n(?=\w)", "", s)
    # Collapse runs of whitespace while preserving single newlines
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\s+\n", "\n", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def _iter_nonempty(pairs: Iterable[Tuple[str, str]]) -> Iterator[Tuple[str, str]]:
    for src, txt in pairs:
        cleaned = _clean_text(txt)
        if cleaned:
            yield (src, cleaned)


def _decode_bytes(data: bytes) -> str:
    """Attempt a best-effort decode without introducing external deps."""
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


# ------------------------------
# HTML
# ------------------------------

def iter_html(path: str) -> Iterator[Tuple[str, str]]:
    """Extract HTML with lightweight structure.

    Emits segments labeled by section heading order when possible:
    - base#sec=N for content under a heading (h1..h6)
    - Falls back to base for uncategorized content
    """
    with open(path, "rb") as f:
        raw = f.read()
    try:
        soup = BeautifulSoup(raw, "lxml")
    except Exception:
        soup = BeautifulSoup(raw, "html.parser")
    for bad in soup(["script", "style", "noscript", "header", "footer", "nav", "aside"]):
        bad.decompose()

    base = os.path.basename(path)
    sec_index = 0
    el_index = 0
    current_sec = None
    # Traverse in document order
    for element in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "th", "td"]):
        content = element.get_text(" ", strip=True)
        if not content:
            continue
        tag = element.name.lower()
        if tag in {"h1", "h2", "h3", "h4", "h5", "h6"}:
            sec_index += 1
            current_sec = sec_index
            el_index = 0
            yield (f"{base}#sec={current_sec}", content)
        else:
            el_index += 1
            if current_sec is not None:
                yield (f"{base}#sec={current_sec}#el={el_index}", content)
            else:
                yield (base, content)
    # Fallback: if nothing yielded, return whole text
    # (this rarely triggers because the loop above covers most tags)
    # Keep as a guard for very minimal HTMLs


# ------------------------------
# RTF
# ------------------------------

def _iter_rtf_images(raw: bytes, base_name: str) -> Iterator[Tuple[str, str]]:
    text = raw.decode("latin-1", errors="ignore")
    pict_re = re.compile(r"\\pict[^{]*?([0-9A-Fa-f\\s]+)}", re.DOTALL)
    for idx, match in enumerate(pict_re.finditer(text), start=1):
        hex_blob = re.sub(r"[^0-9A-Fa-f]", "", match.group(1))
        if len(hex_blob) < 20:
            continue
        try:
            data = bytes.fromhex(hex_blob)
        except ValueError:
            continue
        ocr = _ocr_image_bytes(data)
        if ocr:
            yield (f"{base_name}#image={idx}", ocr)


def iter_rtf(path: str) -> Iterator[Tuple[str, str]]:
    with open(path, "rb") as f:
        raw = f.read()
    try:
        txt = rtf_to_text(raw.decode("utf-8", errors="ignore"))
    except Exception:
        txt = rtf_to_text(raw.decode("latin-1", errors="ignore"))
    base = os.path.basename(path)
    yield (base, txt)
    yield from _iter_rtf_images(raw, base)


# ------------------------------
# DOCX
# ------------------------------

def _iter_docx_images(path: str) -> Iterator[Tuple[str, str]]:
    base = os.path.basename(path)
    try:
        with zipfile.ZipFile(path) as zf:
            names = sorted(n for n in zf.namelist() if n.startswith("word/media/"))
            for idx, name in enumerate(names, start=1):
                try:
                    data = zf.read(name)
                except KeyError:
                    continue
                ocr = _ocr_image_bytes(data)
                if ocr:
                    yield (f"{base}#image={idx}", ocr)
    except zipfile.BadZipFile:
        return


def iter_docx(path: str) -> Iterator[Tuple[str, str]]:
    """Extract DOCX as structured segments with stable refs.

    - Paragraphs are emitted individually as base#para=N
    - Table rows are emitted as base#table=T#row=R
    - Images are OCR'd via _iter_docx_images
    """
    doc = DocxDocument(path)
    base = os.path.basename(path)

    # Paragraphs
    para_index = 0
    for para in doc.paragraphs:
        txt = (para.text or "").strip()
        if not txt:
            continue
        para_index += 1
        yield (f"{base}#para={para_index}", txt)

    # Tables
    table_index = 0
    for table in doc.tables:
        table_index += 1
        row_index = 0
        for row in table.rows:
            row_index += 1
            cells = [cell.text.strip() for cell in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                yield (f"{base}#table={table_index}#row={row_index}", line)

    # Images (OCR)
    yield from _iter_docx_images(path)


# ------------------------------
# PPTX
# ------------------------------

def _iter_pptx_images(path: str) -> Iterator[Tuple[str, str]]:
    base = os.path.basename(path)
    try:
        with zipfile.ZipFile(path) as zf:
            names = sorted(n for n in zf.namelist() if n.startswith("ppt/media/"))
            for idx, name in enumerate(names, start=1):
                try:
                    data = zf.read(name)
                except KeyError:
                    continue
                ocr = _ocr_image_bytes(data)
                if ocr:
                    yield (f"{base}#image={idx}", ocr)
    except zipfile.BadZipFile:
        return


def iter_pptx(path: str) -> Iterator[Tuple[str, str]]:
    prs = Presentation(path)
    base = os.path.basename(path)
    for slide_index, slide in enumerate(prs.slides, start=1):
        lines: List[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        lines.append(text)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and hasattr(shape, "image"):
                ocr = _ocr_image_bytes(shape.image.blob)
                if ocr:
                    lines.append(ocr)
        if lines:
            yield (f"{base}#slide={slide_index}", "\n".join(lines))
    yield from _iter_pptx_images(path)


# ------------------------------
# PDF (text-first; OCR fallback)
# ------------------------------

def _try_ocr(img: Image.Image) -> str:
    try:
        import pytesseract
    except Exception:
        return ""
    try:
        langs = os.getenv("TESS_LANGS", "spa")
        return pytesseract.image_to_string(img, lang=langs)
    except Exception:
        try:
            return pytesseract.image_to_string(img)
        except Exception:
            return ""


def _ocr_image_bytes(data: bytes) -> str:
    try:
        with Image.open(io.BytesIO(data)) as raw:
            img = raw.convert("L")
            try:
                return _try_ocr(img)
            finally:
                img.close()
    except Exception:
        return ""


def _extract_pdf_page_text(page: fitz.Page) -> str:
    blocks = page.get_text("blocks")
    if blocks:
        blocks = sorted(blocks, key=lambda b: (round(b[1], 2), round(b[0], 2)))
        lines = [block[4].strip() for block in blocks if len(block) > 4 and block[4].strip()]
        if lines:
            return "\n".join(lines)
    return page.get_text("text") or ""


def iter_pdf(path: str) -> Iterator[Tuple[str, str]]:
    try:
        doc = fitz.open(path)
    except Exception:
        return
    base = os.path.basename(path)
    try:
        for page_index, page in enumerate(doc, start=1):
            text = _extract_pdf_page_text(page)
            if text.strip():
                yield (f"{base}#page={page_index}", text)
                continue
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            png_bytes = pix.tobytes("png")
            try:
                with Image.open(io.BytesIO(png_bytes)) as raw:
                    img = raw.convert("L")
                    try:
                        ocr = _try_ocr(img)
                    finally:
                        img.close()
            except Exception:
                ocr = ""
            if ocr.strip():
                yield (f"{base}#page={page_index}", ocr)
    finally:
        doc.close()


# ------------------------------
# Router
# ------------------------------

def parse_file(path: str, content_type: str) -> Iterator[Tuple[str, str]]:
    ct = (content_type or "").lower()
    name = os.path.basename(path).lower()
    if name.endswith((".html", ".htm")) or "html" in ct:
        yield from _iter_nonempty(iter_html(path))
    elif name.endswith(".rtf") or "rtf" in ct:
        yield from _iter_nonempty(iter_rtf(path))
    elif name.endswith(".docx") or "wordprocessingml.document" in ct:
        yield from _iter_nonempty(iter_docx(path))
    elif name.endswith(".pptx") or "presentationml.presentation" in ct:
        yield from _iter_nonempty(iter_pptx(path))
    elif name.endswith(".pdf") or "pdf" in ct:
        yield from _iter_nonempty(iter_pdf(path))
    else:
        with open(path, "rb") as f:
            data = _decode_bytes(f.read())
        yield from _iter_nonempty([(os.path.basename(path), data)])
